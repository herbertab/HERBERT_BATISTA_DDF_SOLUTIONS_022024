import streamlit as st
from io import BytesIO
import pandas as pd
import os
import requests
import json
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches

# Configurando a chave de API do OpenAI
openai_api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=openai_api_key)

# Carregando dados do arquivo CSV
df = pd.read_csv("corpus-simple_head.csv")
df = df.sort_values(by="Product Name")

# Adiciona uma imagem na apresentação
def add_image(slide, image_path):    
    # Adiciona a nova imagem
    left = Inches(1)
    top = Inches(1)
    slide.shapes.add_picture(image_path, left, top, width=Inches(2))    
    
# Acessa cliente GPT para gerar e armazenar uma imagem
def generate_image(prompt, out_img_path):    
 
    # Chamada à API para gerar a imagem
    response = client.images.generate(
        model="dall-e-2",
        prompt=prompt,
        size="512x512",
        quality="standard",
        n=1,
        response_format='url'
    )

    # Obtendo a URL da imagem gerada
    image_url = response.data[0].url
    
    # Salvando imagem no ambiente
    image_data = requests.get(image_url).content
    with open(out_img_path, 'wb') as f:
        f.write(image_data)

        
# Acessa a API GPT para gerar informação de texto
def generate_text(prompt_user, prompt_system):
    # Fazer a chamada à API GPT
    response = client.chat.completions.create(
        model="gpt-3.5-turbo-0125",  # Modelo GPT-3 
        response_format={ "type": "json_object" },
        messages=[
            {"role": "system", "content": prompt_system},
            {"role": "user", "content": prompt_user}  
        ]        
    )

    # Extrair e retornar a resposta gerada
    return response.choices[0].message.content   
    
        
# Substitui textos na apresentação
def replace_text(replacements, slide):
    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text


# Função principal do Streamlit
def main():
    st.title("Criação de Apresentação PPTX para um produto")

    # Carrega a apresentação
    presentation_path = "template.pptx"
    presentation = Presentation(presentation_path)

    # Escolhe um slide específico (por exemplo, o primeiro slide, índice 0)
    slide_index = 0
    slide = presentation.slides[slide_index]
    
    # Selecionando um produto usando um menu suspenso
    products = df["Product Name"].tolist()
    selected_product = st.selectbox("Select a product:", products)

    # Adiciona um botão para acionar a substituição da imagem
    if st.button("Gerar a Apresentação"):
        
        # Gerar as informações sobre o produto
        # 1ª Requisição: Título da Apresentação
        prompt_system = """
            You are a skilled Power Point presentation creator. You will be provided with a description of a product and must provide the following information in JSON format. 
            STEP 1 - A Title for the power point presentation with a maximum of 50 characters. The title must not have Presentation in it. 
            STEP 2 - A text of up to 300 characters that describes the main characteristics of this product. 
            STEP 3 - Another text of up to 500 characters that describes potential customers, price and advantages in relation to competitors for the product.
            Please observe:
            The outputs must be named as Title, Features and Customer, respectively.
            No other answer should be delivered.
        """        
        target_product = df[df['Product Name'] == selected_product]
        prompt_user = f"Product Description: {target_product['text'].values[0]}"
        
        response = json.loads(generate_text(prompt_user, prompt_system))
        
        replacements = {
            '{T}': response['Title'],
            '{F}': response['Features'],
            '{C}': response['Customer']
        }
        
        # Substituir textos no template pelas informações geradas
        replace_text(replacements, slide)    
        
        
        
        # Gerar a imagem do produto selecionado
        prompt = f"An image of the product: {selected_product}. Only the product in a white background. Nothing else."
        out_img_path = "generated_image.jpg"
        generate_image(prompt, out_img_path)
        
        # Insere a imagem no slide
        add_image(slide, out_img_path)
        
        # Salva a apresentação com a imagem substituída
        output_pptx = BytesIO()
        presentation.save(output_pptx)

        # Oferece o download do arquivo modificado
        st.download_button(
            label="Baixar Apresentação",
            data=output_pptx.getvalue(),
            file_name=f"presentation_{selected_product}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

# Executa o aplicativo Streamlit
if __name__ == "__main__":
    main()
    
